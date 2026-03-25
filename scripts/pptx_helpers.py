# -*- coding: utf-8 -*-
"""
pptx_helpers — 可复用的 PPTX 精细构建元素库

从 4 个实战 PPTX 项目中提炼的元素级 helper 函数。
与 pptx_template.py（整页级模板）互补使用。

使用方式:
    from pptx_helpers import *
    prs, C = create_prs('navy_teal')  # 或 'tech_dark', 'corporate'

依赖:
    pip install python-pptx lxml
"""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ═══════════════════════════════════════════════════
# 配色方案
# ═══════════════════════════════════════════════════

NAVY_TEAL = {
    'navy':        RGBColor(0x1B, 0x3A, 0x5C),
    'navy_dark':   RGBColor(0x12, 0x27, 0x3D),
    'teal':        RGBColor(0x1A, 0x8A, 0x8A),
    'amber':       RGBColor(0xD4, 0x88, 0x0F),
    'red':         RGBColor(0xC0, 0x39, 0x2B),
    'green':       RGBColor(0x27, 0xAE, 0x60),
    'orange':      RGBColor(0xE6, 0x7E, 0x22),
    'white':       RGBColor(0xFF, 0xFF, 0xFF),
    'black':       RGBColor(0x2D, 0x34, 0x36),
    'gray_light':  RGBColor(0xF0, 0xF3, 0xF7),
    'gray_mid':    RGBColor(0xBD, 0xC3, 0xC7),
    'gray_text':   RGBColor(0x5D, 0x6D, 0x7E),
    'row_alt':     RGBColor(0xE8, 0xED, 0xF2),
    'light_blue':  RGBColor(0xD6, 0xE4, 0xF0),
    'light_red':   RGBColor(0xFD, 0xED, 0xED),
    'light_green': RGBColor(0xE8, 0xF8, 0xF0),
    'light_amber': RGBColor(0xFE, 0xF5, 0xE7),
    'code_bg':     RGBColor(0x2D, 0x2D, 0x2D),
    'code_text':   RGBColor(0xE0, 0xE0, 0xE0),
    'font_cn':     '\u5fae\u8f6f\u96c5\u9ed1',
    'font_code':   'Courier New',
}

TECH_DARK = {
    'navy':        RGBColor(0x0F, 0x0F, 0x23),
    'navy_dark':   RGBColor(0x0A, 0x0A, 0x1A),
    'teal':        RGBColor(0x00, 0xD9, 0xFF),
    'amber':       RGBColor(0xF3, 0x9C, 0x12),
    'red':         RGBColor(0xE7, 0x4C, 0x3C),
    'green':       RGBColor(0x2E, 0xCC, 0x71),
    'orange':      RGBColor(0xE6, 0x7E, 0x22),
    'white':       RGBColor(0xFF, 0xFF, 0xFF),
    'black':       RGBColor(0xE0, 0xE0, 0xE0),  # light text on dark bg
    'gray_light':  RGBColor(0x1A, 0x1A, 0x3E),
    'gray_mid':    RGBColor(0xAA, 0xAA, 0xAA),
    'gray_text':   RGBColor(0x88, 0x88, 0xAA),
    'row_alt':     RGBColor(0x1F, 0x1F, 0x40),
    'light_blue':  RGBColor(0x1A, 0x1A, 0x3E),
    'light_red':   RGBColor(0x3E, 0x1A, 0x1A),
    'light_green': RGBColor(0x1A, 0x3E, 0x2A),
    'light_amber': RGBColor(0x3E, 0x30, 0x1A),
    'code_bg':     RGBColor(0x1E, 0x1E, 0x1E),
    'code_text':   RGBColor(0xD4, 0xD4, 0xD4),
    'font_cn':     '\u5fae\u8f6f\u96c5\u9ed1',
    'font_code':   'Consolas',
}

CORPORATE = {
    'navy':        RGBColor(0x2E, 0x5C, 0x9A),
    'navy_dark':   RGBColor(0x1C, 0x3F, 0x6E),
    'teal':        RGBColor(0x2E, 0x5C, 0x9A),
    'amber':       RGBColor(0xD4, 0x88, 0x0F),
    'red':         RGBColor(0xC0, 0x00, 0x00),
    'green':       RGBColor(0x27, 0xAE, 0x60),
    'orange':      RGBColor(0xE6, 0x7E, 0x22),
    'white':       RGBColor(0xFF, 0xFF, 0xFF),
    'black':       RGBColor(0x00, 0x00, 0x00),
    'gray_light':  RGBColor(0xF5, 0xF5, 0xF5),
    'gray_mid':    RGBColor(0xD9, 0xD9, 0xD9),
    'gray_text':   RGBColor(0x66, 0x66, 0x66),
    'row_alt':     RGBColor(0xF0, 0xF0, 0xF0),
    'light_blue':  RGBColor(0xDD, 0xEB, 0xF7),
    'light_red':   RGBColor(0xFD, 0xED, 0xED),
    'light_green': RGBColor(0xE2, 0xEF, 0xDA),
    'light_amber': RGBColor(0xFE, 0xF5, 0xE7),
    'code_bg':     RGBColor(0xF5, 0xF5, 0xF5),
    'code_text':   RGBColor(0x33, 0x33, 0x33),
    'font_cn':     '\u5fae\u8f6f\u96c5\u9ed1',
    'font_code':   'Courier New',
}

SCHEMES = {
    'navy_teal': NAVY_TEAL,
    'tech_dark': TECH_DARK,
    'corporate': CORPORATE,
}

# Default slide dimensions (16:9 EMU)
SLIDE_W = 12192000
SLIDE_H = 6858000


# ═══════════════════════════════════════════════════
# 工厂函数
# ═══════════════════════════════════════════════════

def create_prs(scheme='navy_teal'):
    """创建 Presentation 并返回 (prs, color_dict)。

    Args:
        scheme: 'navy_teal' | 'tech_dark' | 'corporate' 或自定义 dict

    Returns:
        (Presentation, dict) — prs 和配色方案字典
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    C = SCHEMES[scheme] if isinstance(scheme, str) else scheme
    return prs, C


# ═══════════════════════════════════════════════════
# 元素级 Helper 函数
# ═══════════════════════════════════════════════════

def add_bg(s, l, t, w, h, c):
    """添加纯色矩形背景块。

    Args:
        s: slide  l,t,w,h: left/top/width/height (EMU)  c: RGBColor
    Returns:
        shape
    """
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def tb(s, l, t, w, h, txt, sz=14, b=False, c=None, al=PP_ALIGN.LEFT,
       fn=None, ls=1.15, C=None):
    """添加单行/多行文本框。

    Args:
        s: slide  txt: 文本  sz: 字号  b: 粗体  c: 颜色
        al: 对齐  fn: 字体  ls: 行距  C: 配色方案
    """
    if C is None:
        C = NAVY_TEAL
    if c is None:
        c = C['black']
    if fn is None:
        fn = C['font_cn']
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = al
    p.line_spacing = ls
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.name = fn
    r.font.color.rgb = c
    return bx


def ml(s, l, t, w, h, lines, sz=12, c=None, b=False, al=PP_ALIGN.LEFT,
       bullet=False, ls=1.3, fn=None, C=None):
    """添加多行文本框。

    Args:
        lines: str 列表  bullet: 是否添加项目符号
    """
    if C is None:
        C = NAVY_TEAL
    if c is None:
        c = C['black']
    if fn is None:
        fn = C['font_cn']
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        p.line_spacing = ls
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        if bullet:
            pPr = p._p.get_or_add_pPr()
            bc = etree.SubElement(pPr, qn('a:buChar'))
            bc.set('char', '\u2022')
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.name = fn
        r.font.color.rgb = c
    return bx


def rl(s, l, t, w, h, items, ls=1.3, C=None):
    """添加富文本列表（每行独立样式）。

    Args:
        items: [(text, size, bold, color), ...]
    """
    if C is None:
        C = NAVY_TEAL
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bd, cl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = ls
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        pPr = p._p.get_or_add_pPr()
        bc = etree.SubElement(pPr, qn('a:buChar'))
        bc.set('char', '\u2022')
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz or 12)
        r.font.bold = bd
        r.font.name = C['font_cn']
        r.font.color.rgb = cl or C['black']
    return bx


def sc(cell, txt, sz=10, b=False, fc=None, bg=None, al=PP_ALIGN.LEFT,
       an=MSO_ANCHOR.MIDDLE, fn=None, C=None):
    """设置表格单元格内容和样式。

    Args:
        cell: table.cell(r,c)  txt: 文本  fc: 字体色  bg: 背景色
    """
    if C is None:
        C = NAVY_TEAL
    if fc is None:
        fc = C['black']
    if fn is None:
        fn = C['font_cn']
    tf = cell.text_frame
    tf.word_wrap = True
    try:
        cell.vertical_anchor = an
    except Exception:
        pass
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf.paragraphs[0].text = ''
    tf.margin_left = Emu(72000)
    tf.margin_right = Emu(72000)
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)
    for i, line in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.name = fn
        r.font.color.rgb = fc


def ct(s, l, t, w, h, rows, cols, cw=None):
    """创建表格并返回 table 对象。

    Args:
        cw: 列宽列表 (EMU)
    """
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    _tbl = tbl._tbl
    pr = _tbl.find(qn('a:tblPr'))
    if pr is None:
        pr = etree.SubElement(_tbl, qn('a:tblPr'))
    pr.set('firstRow', '1')
    pr.set('bandRow', '1')
    if cw:
        for i, v in enumerate(cw):
            tbl.columns[i].width = v
    return tbl


def ft(tbl, hdr, data, hc=None, hfc=None, fs=10, hfs=11, C=None):
    """填充表格（表头 + 数据行，自动交替行色）。

    Args:
        hdr: 表头列表  data: 二维数据列表
        hc: 表头背景色  hfc: 表头字体色
    """
    if C is None:
        C = NAVY_TEAL
    if hc is None:
        hc = C['navy']
    if hfc is None:
        hfc = C['white']
    for j, h in enumerate(hdr):
        sc(tbl.cell(0, j), h, hfs, True, hfc, hc, PP_ALIGN.CENTER, C=C)
    for i, row in enumerate(data):
        bg_r = C['row_alt'] if i % 2 == 0 else None
        for j, v in enumerate(row):
            sc(tbl.cell(i + 1, j), v, fs, False, C['black'], bg_r, C=C)


def code(s, l, t, w, h, lines, sz=10, C=None):
    """添加深色背景代码块。"""
    if C is None:
        C = NAVY_TEAL
    add_bg(s, l, t, w, h, C['code_bg'])
    bx = s.shapes.add_textbox(l + Emu(60000), t + Emu(40000),
                               w - Emu(120000), h - Emu(80000))
    tf = bx.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        p.line_spacing = 1.2
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(sz)
        r.font.name = C['font_code']
        r.font.color.rgb = C['code_text']


def hbox(s, l, t, w, h, title, lines, bg=None, tc=None, C=None):
    """添加带标题的高亮信息框。"""
    if C is None:
        C = NAVY_TEAL
    if bg is None:
        bg = C['light_blue']
    if tc is None:
        tc = C['navy']
    add_bg(s, l, t, w, h, bg)
    tb(s, l + Emu(80000), t + Emu(60000), w - Emu(160000), Emu(280000),
       title, sz=13, b=True, c=tc, C=C)
    if lines:
        ml(s, l + Emu(80000), t + Emu(330000), w - Emu(160000), h - Emu(380000),
           lines, sz=11, c=C['black'], ls=1.2, bullet=True, C=C)


def badge(s, l, t, w, h, txt, bg=None, C=None):
    """添加序号/标签徽章。"""
    if C is None:
        C = NAVY_TEAL
    if bg is None:
        bg = C['teal']
    sh = add_bg(s, l, t, w, h, bg)
    tf = sh.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.name = C['font_cn']
    r.font.color.rgb = C['white']
    return sh


def sn(s, n, total, C=None):
    """添加页码标注 (右下角)。"""
    if C is None:
        C = NAVY_TEAL
    tb(s, Emu(SLIDE_W - 900000), Emu(SLIDE_H - 350000),
       Emu(700000), Emu(250000),
       f'{n} / {total}', sz=9, c=C['gray_text'], al=PP_ALIGN.RIGHT, C=C)


def bar(s, title, sub=None, C=None):
    """添加顶部装饰条 + 标题 + 底部条。"""
    if C is None:
        C = NAVY_TEAL
    add_bg(s, 0, 0, SLIDE_W, Emu(75000), C['navy'])
    add_bg(s, 0, Emu(75000), SLIDE_W, Emu(32000), C['teal'])
    tb(s, Emu(450000), Emu(180000), Emu(10000000), Emu(450000),
       title, sz=24, b=True, c=C['navy'], C=C)
    if sub:
        tb(s, Emu(450000), Emu(580000), Emu(10000000), Emu(300000),
           sub, sz=13, c=C['gray_text'], C=C)
    add_bg(s, 0, Emu(SLIDE_H - 100000), SLIDE_W, Emu(100000), C['navy'])


def act_badge(s, l, t, act_num, act_title, act_color, C=None):
    """添加章节标记条（幕标题）。"""
    if C is None:
        C = NAVY_TEAL
    add_bg(s, l, t, Emu(11300000), Emu(45000), act_color)
    badge(s, l, t + Emu(60000), Emu(1200000), Emu(340000),
          f'\u7b2c{act_num}\u5e55', act_color, C=C)
    tb(s, l + Emu(1400000), t + Emu(75000), Emu(8000000), Emu(300000),
       act_title, sz=15, b=True, c=act_color, C=C)


# ═══════════════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════════════

def save_pptx(prs, path):
    """保存并打印信息。"""
    prs.save(path)
    import os
    size_kb = os.path.getsize(path) / 1024
    print(f'Saved: {path}')
    print(f'Total slides: {len(prs.slides)}')
    print(f'Size: {size_kb:.0f} KB')


def sanitize_script(path):
    """检查并修复 Python 脚本中的 Unicode 引号问题。

    返回 True 如果已修复或无问题，False 如果有其他语法错误。
    """
    import ast
    with open(path, 'r', encoding='utf-8') as f:
        content = f.read()
    # 检查并替换 Unicode 引号
    replacements = {
        '\u201c': '"',   # 左双引号
        '\u201d': '"',   # 右双引号
        '\u2018': "'",   # 左单引号
        '\u2019': "'",   # 右单引号
    }
    modified = False
    for old, new in replacements.items():
        if old in content:
            content = content.replace(old, new)
            modified = True
    if modified:
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)
        print(f'Fixed Unicode quotes in {path}')
    # 验证语法
    try:
        ast.parse(content)
        return True
    except SyntaxError as e:
        print(f'SyntaxError in {path}: {e}')
        return False
