# -*- coding: utf-8 -*-
"""
PPTX Creator - 通用演示文稿生成模板
使用 python-pptx 生成专业的 PowerPoint 演示文稿

使用方法:
    python pptx_template.py

依赖:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


# ============================================================
# 配色方案 - 可根据主题修改
# ============================================================

# 科技主题配色（深色）
COLORS_TECH = {
    'bg_dark': RGBColor(0x0F, 0x0F, 0x23),
    'accent_blue': RGBColor(0x66, 0x7E, 0xEA),
    'accent_purple': RGBColor(0x76, 0x4B, 0xA2),
    'accent_cyan': RGBColor(0x00, 0xD9, 0xFF),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0xAA, 0xAA, 0xAA),
    'card_bg': RGBColor(0x1A, 0x1A, 0x3E),
    'line_gray': RGBColor(0x33, 0x33, 0x55),
    'font_cn': '微软雅黑',
    'font_en': 'Arial',
    'font_title_cn': '微软雅黑',
}

# 商务主题配色（深色）
COLORS_BUSINESS = {
    'bg_dark': RGBColor(0x1C, 0x28, 0x33),
    'accent_blue': RGBColor(0x2E, 0x86, 0xAB),
    'accent_purple': RGBColor(0x2C, 0x3E, 0x50),
    'accent_cyan': RGBColor(0xF3, 0x9C, 0x12),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0xBD, 0xC3, 0xC7),
    'card_bg': RGBColor(0x2E, 0x40, 0x53),
    'line_gray': RGBColor(0x34, 0x49, 0x5E),
    'font_cn': '微软雅黑',
    'font_en': 'Arial',
    'font_title_cn': '微软雅黑',
}

# 企业主题配色（浅色 - 白底蓝色表头）
COLORS_CORPORATE = {
    'bg_dark': RGBColor(0xFF, 0xFF, 0xFF),   # 白色背景
    'accent_blue': RGBColor(0x2E, 0x5C, 0x9A),  # 蓝色表头
    'accent_purple': RGBColor(0x34, 0x49, 0x5E),
    'accent_cyan': RGBColor(0x2E, 0x5C, 0x9A),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x00, 0x00, 0x00),
    'gray': RGBColor(0x66, 0x66, 0x66),
    'card_bg': RGBColor(0xF5, 0xF5, 0xF5),
    'line_gray': RGBColor(0xD9, 0xD9, 0xD9),
    'accent_red': RGBColor(0xC0, 0x00, 0x00),
    'header_blue': RGBColor(0x2E, 0x5C, 0x9A),
    'text_black': RGBColor(0x00, 0x00, 0x00),
    'text_white': RGBColor(0xFF, 0xFF, 0xFF),
    'border_gray': RGBColor(0xD9, 0xD9, 0xD9),
    'font_cn': '微软雅黑',
    'font_en': 'Arial',
    'font_title_cn': '等线',
}

# 选择配色方案
COLORS = COLORS_TECH


# ============================================================
# 工具函数
# ============================================================

def set_shape_fill(shape, color):
    """设置形状填充颜色"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_text_frame(shape, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    """添加文本到形状

    Args:
        font_name: 字体名称，默认从当前 COLORS 方案中获取 font_cn（中文）或 font_en（英文）
    """
    if font_name is None:
        font_name = COLORS.get('font_cn', 'Arial')
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_slide_background(slide, prs, color):
    """添加幻灯片背景"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    set_shape_fill(background, color)
    background.line.fill.background()
    return background


def add_header_bar(slide, prs, title, color):
    """添加标题栏"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    set_shape_fill(header, color)
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.5))
    add_text_frame(title_box, title, 28, COLORS['white'], bold=True)
    return header


# ============================================================
# 表格工具函数
# ============================================================

def create_table(slide, left, top, width, height, rows, cols, col_widths=None):
    """创建表格并返回 table 对象

    Args:
        slide: 幻灯片对象
        left, top, width, height: 表格位置和大小（EMU 或 Inches/Pt 值）
        rows: 行数
        cols: 列数
        col_widths: 可选，各列宽度列表（EMU 或 Inches 值）

    Returns:
        table 对象
    """
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            if i < len(table.columns):
                table.columns[i].width = w

    return table


def set_table_style(table, style_id='{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}',
                    first_row=True, first_col=False, band_row=True):
    """设置表格样式（通过 lxml 操作 XML，python-pptx 没有直接 API）

    Args:
        table: python-pptx table 对象
        style_id: 表格样式 GUID，默认 Medium Style 2 - Accent 1
        first_row: 是否启用首行格式
        first_col: 是否启用首列格式
        band_row: 是否启用交替行
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    tbl = table._tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = etree.SubElement(tbl, qn('a:tblPr'))

    tblPr.set('firstRow', '1' if first_row else '0')
    tblPr.set('firstCol', '1' if first_col else '0')
    tblPr.set('bandRow', '1' if band_row else '0')

    # 设置 tableStyleId
    existing = tblPr.find(qn('a:tableStyleId'))
    if existing is not None:
        existing.text = style_id
    else:
        style_elem = etree.SubElement(tblPr, qn('a:tableStyleId'))
        style_elem.text = style_id


def set_cell(cell, text, font_size=10, bold=False, font_color=None,
             fill_color=None, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name=None):
    """设置单个单元格内容和样式

    Args:
        cell: table.cell(r, c) 对象
        text: 文本内容（支持多行，用 \\n 分隔）
        font_size: 字体大小（pt）
        bold: 是否粗体
        font_color: 字体颜色（RGBColor），默认黑色
        fill_color: 填充颜色（RGBColor），None 表示不设置
        alignment: 水平对齐（PP_ALIGN）
        anchor: 垂直锚点（MSO_ANCHOR）
        font_name: 字体名称，默认从 COLORS 方案获取
    """
    if font_color is None:
        font_color = RGBColor(0x00, 0x00, 0x00)
    if font_name is None:
        font_name = COLORS.get('font_cn', 'Arial')

    # 填充色
    if fill_color is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

    # 垂直锚点
    cell.vertical_anchor = anchor

    # 处理多行文本：用 add_paragraph + add_run
    tf = cell.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = font_color
        run.font.name = font_name


def merge_cells(table, r1, c1, r2, c2):
    """合并单元格区域

    注意：必须先调用 merge()，再写入内容。反过来会导致内容丢失。

    Args:
        table: 表格对象
        r1, c1: 起始行列（0-based）
        r2, c2: 结束行列（0-based，包含）
    """
    start_cell = table.cell(r1, c1)
    end_cell = table.cell(r2, c2)
    start_cell.merge(end_cell)


def create_table_slide(prs, title, headers, rows, col_widths=None,
                       header_bg=None, header_font_color=None,
                       font_size=10, header_font_size=11,
                       style_id='{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}'):
    """创建包含表格的幻灯片（高层封装）

    Args:
        prs: Presentation 对象
        title: 幻灯片标题
        headers: 表头列表 ['列1', '列2', ...]
        rows: 数据行列表 [['r1c1', 'r1c2', ...], ['r2c1', ...], ...]
        col_widths: 可选的列宽列表
        header_bg: 表头背景色，默认 accent_blue
        header_font_color: 表头字体色，默认白色
        font_size: 正文字体大小
        header_font_size: 表头字体大小
        style_id: 表格样式 ID

    Returns:
        (slide, table) 元组
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    add_slide_background(slide, prs, COLORS['bg_dark'])

    # 标题栏
    add_header_bar(slide, prs, title, COLORS['accent_blue'])

    # 默认颜色
    if header_bg is None:
        header_bg = COLORS.get('header_blue', COLORS['accent_blue'])
    if header_font_color is None:
        header_font_color = COLORS['white']

    # 计算表格区域
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    table_left = Inches(0.5)
    table_top = Inches(1.1)
    table_width = prs.slide_width - Inches(1.0)
    table_height = prs.slide_height - Inches(1.5)

    table = create_table(slide, table_left, table_top, table_width, table_height,
                         num_rows, num_cols, col_widths)

    # 设置表格样式
    set_table_style(table, style_id=style_id)

    # 填充表头
    for c, header_text in enumerate(headers):
        set_cell(table.cell(0, c), header_text,
                 font_size=header_font_size, bold=True,
                 font_color=header_font_color, fill_color=header_bg,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 填充数据行
    text_color = COLORS.get('text_black', COLORS.get('black', RGBColor(0x00, 0x00, 0x00)))
    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            if c < num_cols:
                set_cell(table.cell(r + 1, c), str(cell_text),
                         font_size=font_size, font_color=text_color)

    return slide, table


# ============================================================
# 模板分析工具
# ============================================================

def analyze_template(pptx_path):
    """分析现有 PPTX 模板结构，输出每个 slide 的 shapes、表格布局、合并情况、样式等

    Args:
        pptx_path: PPTX 文件路径

    Returns:
        分析报告字符串
    """
    prs = Presentation(pptx_path)
    report_lines = []
    report_lines.append(f'=== Template Analysis: {pptx_path} ===')
    report_lines.append(f'Slide size: {prs.slide_width} x {prs.slide_height} EMU '
                        f'({prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} inches)')
    report_lines.append(f'Total slides: {len(prs.slides)}')
    report_lines.append('')

    for i, slide in enumerate(prs.slides):
        report_lines.append(f'--- Slide {i + 1} ---')
        report_lines.append(f'  Layout: {slide.slide_layout.name}')
        report_lines.append(f'  Shapes: {len(slide.shapes)}')

        for j, shape in enumerate(slide.shapes):
            shape_info = f'  [{j}] {shape.shape_type} at ({shape.left}, {shape.top}) size ({shape.width}x{shape.height})'
            report_lines.append(shape_info)

            if shape.has_table:
                t = shape.table
                report_lines.append(f'       Table: {len(t.rows)} rows x {len(t.columns)} cols')

                # 列宽
                col_widths = [col.width for col in t.columns]
                report_lines.append(f'       Col widths: {col_widths}')

                # 检查合并单元格和内容
                for r in range(len(t.rows)):
                    for c in range(len(t.columns)):
                        cell = t.cell(r, c)
                        text = cell.text_frame.text[:60] if cell.text_frame.text else ''
                        if text:
                            report_lines.append(f'       cell({r},{c}): "{text}"')

                # 表格样式
                from pptx.oxml.ns import qn
                tblPr = t._tbl.find(qn('a:tblPr'))
                if tblPr is not None:
                    style_elem = tblPr.find(qn('a:tableStyleId'))
                    if style_elem is not None:
                        report_lines.append(f'       Style ID: {style_elem.text}')

            elif shape.has_text_frame:
                text = shape.text_frame.text[:80]
                if text:
                    report_lines.append(f'       Text: "{text}"')

        report_lines.append('')

    report = '\n'.join(report_lines)
    print(report)
    return report


# ============================================================
# 布局辅助函数
# ============================================================

def add_section_label(slide, left, top, width, height, text,
                      bg_color=None, font_color=None, font_size=14, font_name=None):
    """添加分节标签（如蓝底白字的横条区块标注）

    Args:
        slide: 幻灯片对象
        left, top, width, height: 位置和大小
        text: 标签文字
        bg_color: 背景色，默认 accent_blue
        font_color: 字体色，默认白色
        font_size: 字体大小
        font_name: 字体名称
    """
    if bg_color is None:
        bg_color = COLORS.get('header_blue', COLORS['accent_blue'])
    if font_color is None:
        font_color = COLORS['white']
    if font_name is None:
        font_name = COLORS.get('font_title_cn', COLORS.get('font_cn', 'Arial'))

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_shape_fill(shape, bg_color)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = font_color
    run.font.name = font_name
    p.alignment = PP_ALIGN.LEFT

    return shape


# ============================================================
# 幻灯片创建函数
# ============================================================

def create_title_slide(prs, title, subtitle, tagline=""):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    add_slide_background(slide, prs, COLORS['bg_dark'])
    
    # 顶部强调线
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(8))
    set_shape_fill(top_bar, COLORS['accent_blue'])
    top_bar.line.fill.background()
    
    # 装饰圆形
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(4), Inches(1.2), Inches(1.2))
    set_shape_fill(circle1, COLORS['accent_purple'])
    circle1.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    add_text_frame(title_box, title, 48, COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    add_text_frame(subtitle_box, subtitle, 28, COLORS['accent_blue'], alignment=PP_ALIGN.CENTER)
    
    # 标语
    if tagline:
        tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
        add_text_frame(tagline_box, tagline, 14, COLORS['gray'], alignment=PP_ALIGN.CENTER)
    
    return slide


def create_content_slide(prs, title, items):
    """创建内容列表页
    
    Args:
        prs: Presentation 对象
        title: 幻灯片标题
        items: 列表项 [(标题, 描述), ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    add_slide_background(slide, prs, COLORS['bg_dark'])
    
    # 标题栏
    add_header_bar(slide, prs, title, COLORS['accent_blue'])
    
    # 内容项
    y_start = Inches(1.1)
    for i, (item_title, item_desc) in enumerate(items):
        y = y_start + Inches(i * 0.85)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(9), Inches(0.65)
        )
        set_shape_fill(card, COLORS['card_bg'])
        card.line.fill.background()
        
        # 左边框
        left_border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), y, Pt(4), Inches(0.65)
        )
        set_shape_fill(left_border, COLORS['accent_cyan'])
        left_border.line.fill.background()
        
        # 文本
        text_tb = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.15), Inches(8.5), Inches(0.45))
        tf = text_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = item_title
        run1.font.size = Pt(14)
        run1.font.color.rgb = COLORS['accent_cyan']
        run1.font.bold = True
        run1.font.name = COLORS.get('font_cn', 'Arial')
        if item_desc:
            run2 = p.add_run()
            run2.text = f" - {item_desc}"
            run2.font.size = Pt(14)
            run2.font.color.rgb = COLORS['white']
            run2.font.name = COLORS.get('font_cn', 'Arial')
    
    return slide


def create_cards_slide(prs, title, cards):
    """创建卡片网格页
    
    Args:
        prs: Presentation 对象
        title: 幻灯片标题
        cards: 卡片数据 [(标题, 描述), ...]  最多6个
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    add_slide_background(slide, prs, COLORS['bg_dark'])
    
    # 标题栏
    add_header_bar(slide, prs, title, COLORS['accent_blue'])
    
    # 卡片布局
    card_width = Inches(2.9)
    card_height = Inches(1.8)
    x_positions = [Inches(0.45), Inches(3.55), Inches(6.65)]
    y_positions = [Inches(1.0), Inches(3.0)]
    accent_colors = [COLORS['accent_blue'], COLORS['accent_purple'], COLORS['accent_cyan']]
    
    for i, (card_title, card_desc) in enumerate(cards[:6]):
        col = i % 3
        row = i // 3
        x = x_positions[col]
        y = y_positions[row]
        accent_color = accent_colors[i % 3]
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
        set_shape_fill(card, COLORS['card_bg'])
        card.line.fill.background()
        
        # 顶部边框
        top_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, Pt(4))
        set_shape_fill(top_border, accent_color)
        top_border.line.fill.background()
        
        # 标题
        title_tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), card_width - Inches(0.3), Inches(0.4))
        add_text_frame(title_tb, card_title, 15, COLORS['white'], bold=True)
        
        # 描述
        desc_tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), card_width - Inches(0.3), Inches(1.2))
        tf = desc_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['gray']
        p.font.name = COLORS.get('font_cn', 'Arial')
    
    return slide


def create_summary_slide(prs, title, message):
    """创建总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    add_slide_background(slide, prs, COLORS['bg_dark'])
    
    # 标题栏
    add_header_bar(slide, prs, title, COLORS['accent_blue'])
    
    # 总结框
    summary_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(0.8)
    )
    set_shape_fill(summary_bar, COLORS['accent_blue'])
    summary_bar.line.fill.background()
    
    summary_text = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(8), Inches(0.5))
    add_text_frame(summary_text, message, 18, COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    
    return slide


# ============================================================
# 主函数 - 示例用法
# ============================================================

def main():
    """示例：创建一个演示文稿"""
    
    # 创建演示文稿 (16:9 标准 PowerPoint 默认尺寸)
    prs = Presentation()
    prs.slide_width = 12192000   # 标准 16:9 = Inches(13.333)
    prs.slide_height = 6858000   # 标准 16:9 = Inches(7.5)
    
    print("创建演示文稿...")
    
    # 幻灯片 1: 标题页
    create_title_slide(
        prs,
        title="演示文稿标题",
        subtitle="副标题内容",
        tagline="这是一个使用 python-pptx 生成的演示文稿"
    )
    print("  ✓ 标题页")
    
    # 幻灯片 2: 内容列表
    create_content_slide(
        prs,
        title="主要内容",
        items=[
            ("第一点", "这是第一个要点的详细说明"),
            ("第二点", "这是第二个要点的详细说明"),
            ("第三点", "这是第三个要点的详细说明"),
            ("第四点", "这是第四个要点的详细说明"),
        ]
    )
    print("  ✓ 内容列表页")
    
    # 幻灯片 3: 卡片网格
    create_cards_slide(
        prs,
        title="核心要素",
        cards=[
            ("要素一", "这是第一个要素的详细描述内容，可以包含多行文字。"),
            ("要素二", "这是第二个要素的详细描述内容。"),
            ("要素三", "这是第三个要素的详细描述内容。"),
            ("要素四", "这是第四个要素的详细描述内容。"),
            ("要素五", "这是第五个要素的详细描述内容。"),
            ("要素六", "这是第六个要素的详细描述内容。"),
        ]
    )
    print("  ✓ 卡片网格页")
    
    # 幻灯片 4: 表格页
    create_table_slide(
        prs,
        title="数据概览",
        headers=["项目", "状态", "负责人", "备注"],
        rows=[
            ["任务一", "已完成", "张三", "按时交付"],
            ["任务二", "进行中", "李四", "需要协调"],
            ["任务三", "未开始", "王五", "待排期"],
        ]
    )
    print("  ✓ 表格页")

    # 幻灯片 5: 总结页
    create_summary_slide(
        prs,
        title="总结",
        message="感谢您的关注！"
    )
    print("  ✓ 总结页")
    
    # 保存
    output_path = os.path.join(os.path.dirname(__file__), '..', 'examples', "example-presentation.pptx")
    prs.save(output_path)
    
    print(f"\n✅ 演示文稿创建成功!")
    print(f"📁 输出文件: {output_path}")
    
    return output_path


if __name__ == "__main__":
    main()
