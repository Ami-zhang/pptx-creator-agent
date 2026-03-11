# -*- coding: utf-8 -*-
"""
AI Development Presentation Generator
使用 python-pptx 生成人工智能发展演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
COLORS = {
    'bg_dark': RGBColor(0x0F, 0x0F, 0x23),      # 深蓝黑
    'accent_blue': RGBColor(0x66, 0x7E, 0xEA),   # 电子蓝紫
    'accent_purple': RGBColor(0x76, 0x4B, 0xA2), # 紫色
    'accent_cyan': RGBColor(0x00, 0xD9, 0xFF),   # 青色
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0xAA, 0xAA, 0xAA),
    'card_bg': RGBColor(0x1A, 0x1A, 0x3E),       # 卡片背景
    'line_gray': RGBColor(0x33, 0x33, 0x55),
    'desc_gray': RGBColor(0xBB, 0xBB, 0xBB),
}

def set_shape_fill(shape, color):
    """设置形状填充颜色"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text_frame(shape, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT):
    """添加文本到形状"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = alignment
    return tf

def create_slide1_title(prs):
    """幻灯片1: 标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, COLORS['bg_dark'])
    background.line.fill.background()
    
    # 顶部强调线
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(8))
    set_shape_fill(top_bar, COLORS['accent_blue'])
    top_bar.line.fill.background()
    
    # 装饰圆形
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(4), Inches(1.2), Inches(1.2))
    set_shape_fill(circle1, COLORS['accent_purple'])
    circle1.fill.fore_color.brightness = 0.7  # 透明效果
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.2), Inches(0.6), Inches(0.6))
    set_shape_fill(circle2, COLORS['accent_cyan'])
    circle2.fill.fore_color.brightness = 0.8
    circle2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    add_text_frame(title_box, "人工智能发展", 48, COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    add_text_frame(subtitle_box, "从过去到未来", 28, COLORS['accent_blue'], alignment=PP_ALIGN.CENTER)
    
    # 标语
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
    add_text_frame(tagline_box, "探索AI技术的演进历程与未来趋势", 14, COLORS['gray'], alignment=PP_ALIGN.CENTER)
    
    # 底部线条
    bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.1), Inches(8), Pt(2))
    set_shape_fill(bottom_line, COLORS['line_gray'])
    bottom_line.line.fill.background()

def create_slide2_history(prs):
    """幻灯片2: AI发展历程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, COLORS['bg_dark'])
    background.line.fill.background()
    
    # 标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    set_shape_fill(header, COLORS['accent_blue'])
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.5))
    add_text_frame(title_box, "AI发展历程", 28, COLORS['white'], bold=True)
    
    # 时间线数据
    timeline = [
        ("1956", "达特茅斯会议", "人工智能概念正式诞生"),
        ("1997", "深蓝战胜卡斯帕罗夫", "AI在国际象棋领域超越人类"),
        ("2012", "深度学习突破", "AlexNet引发深度学习革命"),
        ("2016", "AlphaGo战胜李世石", "AI在围棋领域取得里程碑"),
        ("2022+", "生成式AI时代", "ChatGPT开启大语言模型新纪元"),
    ]
    
    y_start = Inches(1.1)
    for i, (year, title, desc) in enumerate(timeline):
        y = y_start + Inches(i * 0.85)
        
        # 年份框
        year_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(0.9), Inches(0.5))
        set_shape_fill(year_box, COLORS['accent_purple'])
        year_box.line.fill.background()
        year_text = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.1), Inches(0.9), Inches(0.35))
        add_text_frame(year_text, year, 13, COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
        
        # 事件框
        event_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.6), y, Inches(7.6), Inches(0.55))
        set_shape_fill(event_box, COLORS['card_bg'])
        event_box.line.fill.background()
        
        # 左边框
        left_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.6), y, Pt(4), Inches(0.55))
        set_shape_fill(left_border, COLORS['accent_cyan'])
        left_border.line.fill.background()
        
        # 事件文本
        event_text = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.12), Inches(7.2), Inches(0.4))
        tf = event_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = title
        run1.font.size = Pt(13)
        run1.font.color.rgb = COLORS['accent_cyan']
        run1.font.bold = True
        run1.font.name = 'Arial'
        run2 = p.add_run()
        run2.text = f" - {desc}"
        run2.font.size = Pt(13)
        run2.font.color.rgb = COLORS['white']
        run2.font.name = 'Arial'

def create_slide3_technologies(prs):
    """幻灯片3: 当前AI核心技术"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, COLORS['bg_dark'])
    background.line.fill.background()
    
    # 标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    set_shape_fill(header, COLORS['accent_blue'])
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.5))
    add_text_frame(title_box, "当前AI核心技术", 28, COLORS['white'], bold=True)
    
    # 技术卡片数据
    techs = [
        ("机器学习", "通过数据训练模型，实现模式识别和预测。包括监督学习、无监督学习和强化学习等方法。", COLORS['accent_blue']),
        ("深度学习", "基于神经网络的多层学习架构，在图像识别、语音处理等领域表现卓越。", COLORS['accent_purple']),
        ("自然语言处理", "使机器理解和生成人类语言。大语言模型(LLM)实现了突破性进展。", COLORS['accent_cyan']),
        ("计算机视觉", "让机器'看懂'图像和视频，应用于人脸识别、自动驾驶、医学影像等领域。", COLORS['accent_blue']),
        ("生成式AI", "能够创造新内容的AI系统，包括文本生成、图像生成、代码生成等能力。", COLORS['accent_purple']),
        ("多模态AI", "整合文本、图像、音频等多种数据类型，实现更全面的智能理解与交互。", COLORS['accent_cyan']),
    ]
    
    card_width = Inches(2.9)
    card_height = Inches(1.8)
    x_positions = [Inches(0.45), Inches(3.55), Inches(6.65)]
    y_positions = [Inches(1.0), Inches(3.0)]
    
    for i, (title, desc, accent_color) in enumerate(techs):
        col = i % 3
        row = i // 3
        x = x_positions[col]
        y = y_positions[row]
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
        set_shape_fill(card, COLORS['card_bg'])
        card.line.fill.background()
        
        # 顶部边框
        top_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, Pt(4))
        set_shape_fill(top_border, accent_color)
        top_border.line.fill.background()
        
        # 标题
        title_tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), card_width - Inches(0.3), Inches(0.4))
        add_text_frame(title_tb, title, 15, COLORS['white'], bold=True)
        
        # 描述
        desc_tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), card_width - Inches(0.3), Inches(1.2))
        tf = desc_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['desc_gray']
        p.font.name = 'Arial'

def create_slide4_applications(prs):
    """幻灯片4: AI应用场景"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, COLORS['bg_dark'])
    background.line.fill.background()
    
    # 标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    set_shape_fill(header, COLORS['accent_blue'])
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.5))
    add_text_frame(title_box, "AI应用场景", 28, COLORS['white'], bold=True)
    
    # 应用场景数据
    apps = [
        ("医疗健康", "辅助诊断、药物研发、健康监测、医学影像分析", COLORS['accent_blue']),
        ("金融服务", "风险评估、智能投顾、反欺诈检测、自动化交易", COLORS['accent_purple']),
        ("智能制造", "质量检测、预测性维护、供应链优化、机器人自动化", COLORS['accent_cyan']),
        ("自动驾驶", "环境感知、路径规划、决策控制、车联网协同", COLORS['accent_cyan']),
        ("教育培训", "个性化学习、智能辅导、自动评估、教育内容生成", COLORS['accent_blue']),
        ("内容创作", "文案写作、图像生成、视频制作、音乐创作", COLORS['accent_purple']),
    ]
    
    item_height = Inches(0.75)
    x_left = Inches(0.4)
    x_right = Inches(5)
    y_start = Inches(1.0)
    
    for i, (title, desc, accent_color) in enumerate(apps):
        col = i // 3
        row = i % 3
        x = x_left if col == 0 else x_right
        y = y_start + row * Inches(1.35)
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.4), item_height)
        set_shape_fill(card, COLORS['card_bg'])
        card.line.fill.background()
        
        # 左边框
        left_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(3), item_height)
        set_shape_fill(left_border, accent_color)
        left_border.line.fill.background()
        
        # 图标占位（圆形）
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.45), Inches(0.45))
        set_shape_fill(icon, COLORS['accent_blue'])
        icon.line.fill.background()
        
        # 标题
        title_tb = slide.shapes.add_textbox(x + Inches(0.7), y + Inches(0.1), Inches(3.5), Inches(0.3))
        add_text_frame(title_tb, title, 13, COLORS['white'], bold=True)
        
        # 描述
        desc_tb = slide.shapes.add_textbox(x + Inches(0.7), y + Inches(0.4), Inches(3.5), Inches(0.35))
        add_text_frame(desc_tb, desc, 9, COLORS['gray'])

def create_slide5_future(prs):
    """幻灯片5: 未来展望"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, COLORS['bg_dark'])
    background.line.fill.background()
    
    # 标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    set_shape_fill(header, COLORS['accent_blue'])
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.5))
    add_text_frame(title_box, "未来展望", 28, COLORS['white'], bold=True)
    
    # 左侧：发展趋势
    trends_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.4))
    tf = trends_title.text_frame
    p = tf.paragraphs[0]
    p.text = "发展趋势"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['accent_cyan']
    p.font.bold = True
    p.font.name = 'Arial'
    
    # 分隔线
    sep_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(4), Pt(2))
    set_shape_fill(sep_line, COLORS['line_gray'])
    sep_line.line.fill.background()
    
    trends = [
        "• 通用人工智能(AGI)持续探索",
        "• AI与量子计算深度融合",
        "• 边缘AI与物联网协同发展",
        "• AI安全与可解释性研究加强",
        "• 多模态大模型能力持续提升",
    ]
    
    y = Inches(1.6)
    for trend in trends:
        trend_tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(4.5), Inches(0.35))
        add_text_frame(trend_tb, trend, 12, COLORS['white'])
        y += Inches(0.38)
    
    # 右侧：面临挑战
    challenges_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.4))
    tf = challenges_title.text_frame
    p = tf.paragraphs[0]
    p.text = "面临挑战"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['accent_cyan']
    p.font.bold = True
    p.font.name = 'Arial'
    
    # 分隔线
    sep_line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.4), Inches(4), Pt(2))
    set_shape_fill(sep_line2, COLORS['line_gray'])
    sep_line2.line.fill.background()
    
    challenges = [
        ("伦理问题", "隐私保护、算法偏见、决策透明度"),
        ("监管合规", "全球AI治理框架持续完善"),
        ("人才培养", "AI人才缺口与教育体系适配"),
    ]
    
    y = Inches(1.6)
    for title, desc in challenges:
        # 挑战卡片
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), y, Inches(4.3), Inches(0.55))
        set_shape_fill(card, COLORS['card_bg'])
        card.line.fill.background()
        
        # 左边框
        left_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), y, Pt(3), Inches(0.55))
        set_shape_fill(left_border, COLORS['accent_purple'])
        left_border.line.fill.background()
        
        # 文本
        text_tb = slide.shapes.add_textbox(Inches(5.4), y + Inches(0.12), Inches(4), Inches(0.4))
        tf = text_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = title
        run1.font.size = Pt(11)
        run1.font.color.rgb = COLORS['accent_cyan']
        run1.font.bold = True
        run1.font.name = 'Arial'
        run2 = p.add_run()
        run2.text = f" - {desc}"
        run2.font.size = Pt(11)
        run2.font.color.rgb = COLORS['white']
        run2.font.name = 'Arial'
        
        y += Inches(0.7)
    
    # 底部总结
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.55))
    set_shape_fill(bottom_bar, COLORS['accent_blue'])
    bottom_bar.line.fill.background()
    
    bottom_text = slide.shapes.add_textbox(Inches(0.8), Inches(4.72), Inches(8.4), Inches(0.4))
    add_text_frame(bottom_text, "AI将重塑各行各业，人机协作是未来趋势", 14, COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)


def main():
    # 创建演示文稿 (16:9)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    print("创建AI发展演示文稿...")
    
    # 创建幻灯片
    print("  创建幻灯片 1: 标题页")
    create_slide1_title(prs)
    
    print("  创建幻灯片 2: AI发展历程")
    create_slide2_history(prs)
    
    print("  创建幻灯片 3: 核心技术")
    create_slide3_technologies(prs)
    
    print("  创建幻灯片 4: 应用场景")
    create_slide4_applications(prs)
    
    print("  创建幻灯片 5: 未来展望")
    create_slide5_future(prs)
    
    # 保存
    output_path = os.path.join(os.path.dirname(__file__), "AI-Development-Presentation.pptx")
    prs.save(output_path)
    
    print(f"\n演示文稿创建成功!")
    print(f"输出文件: {output_path}")
    
    return output_path


if __name__ == "__main__":
    main()
